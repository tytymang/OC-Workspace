import pandas as pd
from pptx import Presentation
import glob

print("--- EXCEL FILES ---")
excel_files = glob.glob('temp_attachments/*.xlsx')
for file in excel_files:
    try:
        # Load all sheets
        xls = pd.ExcelFile(file)
        for sheet_name in xls.sheet_names:
            print(f"\n[File: {file} | Sheet: {sheet_name}]")
            df = pd.read_excel(file, sheet_name=sheet_name)
            # Find relevant columns or print first few rows to understand structure
            print(df.head(10).to_string())
    except Exception as e:
        print(f"Error reading {file}: {e}")

print("\n--- PPT FILES ---")
ppt_files = glob.glob('temp_attachments/*.pptx')
for file in ppt_files:
    try:
        print(f"\n[File: {file}]")
        prs = Presentation(file)
        for i, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text.replace("\n", " | "))
            if text:
                print(f"Slide {i+1}: {' /// '.join(text)}")
    except Exception as e:
        print(f"Error reading {file}: {e}")
